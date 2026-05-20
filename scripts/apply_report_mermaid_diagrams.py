from __future__ import annotations

import argparse
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


REPO_ROOT = Path(__file__).resolve().parents[1]
SRC_ROOT = REPO_ROOT / "src"
if str(SRC_ROOT) not in sys.path:
    sys.path.insert(0, str(SRC_ROOT))

from openscenario_mcp.reporting.ppt_mermaid import (  # noqa: E402
    asset_path_for_placement,
    extract_mermaid_blocks,
    report_diagram_placements,
)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Export Mermaid sources from the report and replace the target PPT slides with rendered diagram images."
    )
    parser.add_argument("--md", type=Path, required=True)
    parser.add_argument("--ppt", type=Path, required=True)
    parser.add_argument("--output-ppt", type=Path, default=None)
    parser.add_argument("--assets-dir", type=Path, required=True)
    return parser


def export_mermaid_sources(markdown_path: Path, assets_dir: Path) -> list[Path]:
    markdown_text = markdown_path.read_text(encoding="utf-8")
    blocks = extract_mermaid_blocks(markdown_text)
    placements = report_diagram_placements()
    if len(blocks) <= placements[-1].markdown_index:
        raise ValueError(
            f"Markdown contains {len(blocks)} mermaid blocks, but replacements expect at least "
            f"{placements[-1].markdown_index + 1}."
        )

    assets_dir.mkdir(parents=True, exist_ok=True)
    source_paths: list[Path] = []
    for placement in placements:
        block = blocks[placement.markdown_index]
        if placement.expected_fragment not in block:
            raise ValueError(
                f"Mermaid block {placement.markdown_index} does not contain expected fragment "
                f"{placement.expected_fragment!r} for placement {placement.key}."
            )
        source_path = assets_dir / f"{placement.markdown_index + 1:02d}-{placement.key}.mmd"
        source_path.write_text(block + "\n", encoding="utf-8")
        source_paths.append(source_path)
    return source_paths


def apply_diagram_images(source_ppt_path: Path, assets_dir: Path, output_ppt_path: Path | None = None) -> Path:
    target_path = source_ppt_path if output_ppt_path is None else output_ppt_path
    if output_ppt_path is not None:
        output_ppt_path.parent.mkdir(parents=True, exist_ok=True)
        shutil.copyfile(source_ppt_path, output_ppt_path)

    presentation = Presentation(str(target_path))
    for placement in report_diagram_placements():
        image_path = asset_path_for_placement(assets_dir, placement)
        if not image_path.is_file():
            raise FileNotFoundError(f"Missing rendered diagram image: {image_path}")

        slide = presentation.slides[placement.slide_number - 1]
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(placement.left),
            Inches(placement.top),
            Inches(placement.width),
            Inches(placement.height),
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        background.line.fill.background()
        slide.shapes.add_picture(
            str(image_path),
            Inches(placement.left),
            Inches(placement.top),
            Inches(placement.width),
            Inches(placement.height),
        )

    presentation.save(str(target_path))
    return target_path


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    export_mermaid_sources(args.md, args.assets_dir)
    updated_path = apply_diagram_images(args.ppt, args.assets_dir, args.output_ppt)
    print(f"Updated PPT diagrams in {updated_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
