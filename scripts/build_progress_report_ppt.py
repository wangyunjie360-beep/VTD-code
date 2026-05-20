from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


REPORT_DIR = Path("进展汇报")
TEMPLATE_NAME = "项目研究计划(1).pptx"
OUTPUT_NAME = "OpenSCENARIO智能场景生成辅助平台-技术评审汇报.pptx"
BACKGROUND_DIR = Path(".codex-tmp") / "ppt-template-preview"
BACKGROUND_IMAGE = BACKGROUND_DIR / "幻灯片2.PNG"

PURPLE = RGBColor(126, 74, 152)
PURPLE_DARK = RGBColor(97, 53, 122)
BLUE = RGBColor(79, 109, 189)
TEAL = RGBColor(42, 138, 125)
ORANGE = RGBColor(217, 132, 58)
RED = RGBColor(185, 65, 65)
TEXT = RGBColor(35, 42, 58)
TEXT_LIGHT = RGBColor(88, 96, 115)
BORDER = RGBColor(200, 206, 224)
LIGHT = RGBColor(247, 246, 251)
LIGHT_BLUE = RGBColor(238, 244, 255)
LIGHT_TEAL = RGBColor(237, 248, 245)
LIGHT_ORANGE = RGBColor(255, 246, 236)
LIGHT_RED = RGBColor(253, 241, 241)
WHITE = RGBColor(255, 255, 255)

FONT_SANS = "Microsoft YaHei"
FONT_MONO = "Consolas"


def _template_path() -> Path:
    return REPORT_DIR / TEMPLATE_NAME


def _output_path() -> Path:
    return REPORT_DIR / OUTPUT_NAME


def _background_image_path() -> Path:
    if BACKGROUND_IMAGE.exists():
        return BACKGROUND_IMAGE
    import win32com.client  # type: ignore

    BACKGROUND_DIR.mkdir(parents=True, exist_ok=True)
    app = win32com.client.Dispatch("PowerPoint.Application")
    app.Visible = -1
    pres = app.Presentations.Open(str(_template_path().resolve()), False, True, False)
    pres.Export(str(BACKGROUND_DIR.resolve()), "PNG")
    pres.Close()
    app.Quit()
    return BACKGROUND_IMAGE


def _set_font(run, name: str = FONT_SANS, size: int = 18, bold: bool = False, color=TEXT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    r_pr = run._r.get_or_add_rPr()
    r_pr.set(qn("a:ea"), name)
    r_pr.set(qn("a:cs"), name)
    return run


def _add_textbox(slide, left, top, width, height, text="", size=18, bold=False, color=TEXT, name=FONT_SANS, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    _set_font(run, name=name, size=size, bold=bold, color=color)
    return box


def _add_paragraph(frame, text: str, size=16, bold=False, color=TEXT, name=FONT_SANS, align=PP_ALIGN.LEFT, space_before=0):
    paragraph = frame.add_paragraph()
    paragraph.alignment = align
    paragraph.space_before = Pt(space_before)
    run = paragraph.add_run()
    run.text = text
    _set_font(run, name=name, size=size, bold=bold, color=color)
    return paragraph


def _style_shape(shape, fill_color=WHITE, line_color=BORDER, line_width=1.2):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(line_width)
    return shape


def _add_card(slide, left, top, width, height, title, body_lines, fill=WHITE, line=PURPLE, title_color=PURPLE_DARK, body_color=TEXT, title_size=17, body_size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    _style_shape(shape, fill_color=fill, line_color=line, line_width=1.1)
    shape.adjustments[0] = 0.08
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Pt(12)
    frame.margin_right = Pt(12)
    frame.margin_top = Pt(10)
    frame.margin_bottom = Pt(8)
    frame.word_wrap = True
    first = frame.paragraphs[0]
    first.alignment = PP_ALIGN.LEFT
    run = first.add_run()
    run.text = title
    _set_font(run, size=title_size, bold=True, color=title_color)
    for line_text in body_lines:
        _add_paragraph(frame, line_text, size=body_size, color=body_color, space_before=2)
    return shape


def _add_title(slide, title: str, kicker: str):
    _add_textbox(slide, 0.55, 0.44, 7.5, 0.45, title, size=26, bold=True, color=PURPLE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.57), Inches(0.96), Inches(0.95), Inches(0.045))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()
    kicker_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.45), Inches(0.42), Inches(3.0), Inches(0.40))
    _style_shape(kicker_shape, fill_color=LIGHT_BLUE, line_color=BLUE, line_width=1.0)
    kicker_shape.adjustments[0] = 0.12
    frame = kicker_shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = kicker
    _set_font(r, size=12, bold=True, color=BLUE)


def _add_bullets(slide, left, top, width, height, lines, size=15, color=TEXT, fill=None, border=None, title=None):
    if fill is not None:
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        _style_shape(panel, fill_color=fill, line_color=border or BORDER, line_width=1.0)
        panel.adjustments[0] = 0.08
        frame = panel.text_frame
        frame.clear()
        frame.margin_left = Pt(12)
        frame.margin_right = Pt(12)
        frame.margin_top = Pt(10)
        frame.margin_bottom = Pt(8)
    else:
        panel = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = panel.text_frame
        frame.clear()
    frame.word_wrap = True
    if title:
        p = frame.paragraphs[0]
        r = p.add_run()
        r.text = title
        _set_font(r, size=17, bold=True, color=PURPLE_DARK)
    else:
        frame.paragraphs[0].clear()
    for idx, line in enumerate(lines):
        paragraph = frame.add_paragraph()
        paragraph.space_before = Pt(3 if idx else 5)
        run = paragraph.add_run()
        run.text = "• " + line
        _set_font(run, size=size, color=color)
    return panel


def _add_chip(slide, left, top, width, height, text, fill=LIGHT_BLUE, line=BLUE, color=BLUE, size=11):
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    _style_shape(chip, fill_color=fill, line_color=line, line_width=1.0)
    chip.adjustments[0] = 0.25
    tf = chip.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, bold=True, color=color)
    return chip


def _add_code_block(slide, left, top, width, height, lines, title=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    _style_shape(box, fill_color=RGBColor(248, 249, 252), line_color=RGBColor(214, 219, 231), line_width=1.0)
    box.adjustments[0] = 0.05
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Pt(10)
    frame.margin_right = Pt(10)
    frame.margin_top = Pt(10)
    frame.margin_bottom = Pt(6)
    frame.word_wrap = True
    if title:
        p = frame.paragraphs[0]
        r = p.add_run()
        r.text = title
        _set_font(r, size=12, bold=True, color=TEXT_LIGHT)
    else:
        frame.paragraphs[0].clear()
    for line in lines:
        paragraph = frame.add_paragraph()
        run = paragraph.add_run()
        run.text = line
        _set_font(run, name=FONT_MONO, size=11, color=TEXT)
    return box


def _add_metric_card(slide, left, top, width, height, value, label, fill=LIGHT_BLUE, line=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    _style_shape(shape, fill_color=fill, line_color=line, line_width=1.2)
    shape.adjustments[0] = 0.1
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Pt(8)
    frame.margin_right = Pt(8)
    frame.margin_top = Pt(8)
    frame.word_wrap = True
    p1 = frame.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    _set_font(r1, size=24, bold=True, color=PURPLE_DARK)
    p2 = frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    _set_font(r2, size=12, color=TEXT_LIGHT)
    return shape


def _add_flow_row(slide, left, top, width, height, steps, fills=None, text_size=14):
    count = len(steps)
    gap = 0.12
    node_width = (width - gap * (count - 1)) / count
    fills = fills or [LIGHT] * count
    for index, step in enumerate(steps):
        x = left + index * (node_width + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(node_width), Inches(height))
        _style_shape(box, fill_color=fills[index], line_color=PURPLE if index == 0 else BLUE, line_width=1.0)
        box.adjustments[0] = 0.08
        tf = box.text_frame
        tf.clear()
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(6)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = step
        _set_font(r, size=text_size, bold=True, color=TEXT)
        if index < count - 1:
            arrow_left = x + node_width
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(arrow_left + 0.01), Inches(top + height / 2 - 0.16), Inches(gap - 0.02), Inches(0.32))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = PURPLE
            arrow.line.fill.background()


def _add_table(slide, left, top, width, height, data, col_widths, header_fill=PURPLE, row_fill=WHITE):
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    total = sum(col_widths)
    for c, ratio in enumerate(col_widths):
        table.columns[c].width = int(Inches(width) * ratio / total)
    row_height = int(Inches(height) / rows)
    for row in table.rows:
        row.height = row_height
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r == 0 else row_fill
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c == 1 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(data[r][c])
            _set_font(run, size=12 if r == 0 else 11, bold=(r == 0), color=WHITE if r == 0 else TEXT)
            cell.margin_left = Pt(8)
            cell.margin_right = Pt(8)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
    return table


def _add_connector(slide, x1, y1, x2, y2, color=BLUE, width=1.8):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def _clear_all_slides(prs: Presentation):
    for index in range(len(prs.slides) - 1, -1, -1):
        slide_id = prs.slides._sldIdLst[index]
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)


def _new_slide(prs: Presentation, title: str, kicker: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(_background_image_path()), 0, 0, prs.slide_width, prs.slide_height)
    _add_title(slide, title, kicker)
    return slide


def build_cover(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(_background_image_path()), 0, 0, prs.slide_width, prs.slide_height)
    _add_textbox(slide, 0.75, 0.55, 10.5, 0.7, "OpenSCENARIO 智能场景生成辅助平台", size=30, bold=True, color=PURPLE)
    _add_textbox(slide, 0.78, 1.10, 8.2, 0.45, "技术评审汇报", size=20, bold=True, color=BLUE)
    _add_chip(slide, 0.80, 1.60, 1.65, 0.34, "面向 VTD 开发", fill=LIGHT_TEAL, line=TEAL, color=TEAL)
    _add_chip(slide, 2.55, 1.60, 1.55, 0.34, "LLM-first", fill=LIGHT_BLUE, line=BLUE, color=BLUE)
    _add_chip(slide, 4.20, 1.60, 1.95, 0.34, "MCP-assisted", fill=LIGHT_BLUE, line=BLUE, color=BLUE)
    _add_chip(slide, 6.25, 1.60, 3.10, 0.34, "ASAM OpenSCENARIO XML V1.4.0", fill=LIGHT_ORANGE, line=ORANGE, color=ORANGE)
    _add_bullets(
        slide,
        0.78,
        2.10,
        6.0,
        1.9,
        [
            "围绕 OpenSCENARIO XML 生成、修复、校验构建完整技术闭环",
            "采用“结构化知识图谱 + MCP 工具链 + Skill 工作流 + 回归验证”架构",
            "目标不是模板拼接器，而是面向复杂场景的智能辅助平台",
        ],
        size=16,
        fill=WHITE,
        border=BORDER,
        title="评审主结论",
    )
    _add_card(
        slide,
        7.15,
        2.15,
        5.2,
        1.9,
        "本次评审范围",
        [
            "知识图谱底座构建与覆盖状态",
            "MCP 五工具设计与协同方式",
            "Skill / 外部 agent 工作流",
            "Benchmark 与自动化验证结果",
        ],
        fill=LIGHT,
        line=PURPLE,
    )
    _add_metric_card(slide, 0.85, 4.55, 2.7, 1.1, "302", "XSD 元素名覆盖")
    _add_metric_card(slide, 3.75, 4.55, 2.7, 1.1, "301", "结构化 JSON 文件")
    _add_metric_card(slide, 6.65, 4.55, 2.7, 1.1, "5", "核心 MCP 工具")
    _add_metric_card(slide, 9.55, 4.55, 2.7, 1.1, "93 passed", "最新完整回归")
    _add_flow_row(
        slide,
        0.9,
        6.05,
        11.4,
        0.62,
        ["结构化知识", "MCP Runtime", "Skill / Agent", "本地校验闭环", "VTD 导向输出"],
        fills=[LIGHT, LIGHT_BLUE, LIGHT_TEAL, LIGHT_ORANGE, LIGHT_RED],
        text_size=13,
    )


def build_scope_slide(prs: Presentation):
    slide = _new_slide(prs, "评审范围与核心结论", "第 1 部分 / 总览")
    _add_card(
        slide,
        0.65,
        1.35,
        6.05,
        4.9,
        "本次评审重点",
        [
            "1. 知识图谱是否实现“全量覆盖 + 结构可计算”",
            "2. MCP 是否形成从检索到修复的完整工具链",
            "3. Skill 是否把大模型能力与本地结构约束有效衔接",
            "4. Benchmark、回归测试和 sidecar 契约是否完整",
            "5. 当前阶段成果是否足以支撑下一步 VTD 场景深化",
        ],
        fill=WHITE,
        line=PURPLE,
    )
    _add_card(
        slide,
        7.0,
        1.35,
        5.55,
        1.45,
        "结论 01",
        ["已完成基于 XSD 的结构化知识底座，不再依赖零散文档检索。"],
        fill=LIGHT_BLUE,
        line=BLUE,
    )
    _add_card(
        slide,
        7.0,
        2.95,
        5.55,
        1.45,
        "结论 02",
        ["已完成 5 个 MCP 工具及其协同闭环，具备真实校验与结构化诊断能力。"],
        fill=LIGHT_TEAL,
        line=TEAL,
    )
    _add_card(
        slide,
        7.0,
        4.55,
        5.55,
        1.45,
        "结论 03",
        ["已形成 LLM-first、MCP-assisted 的工作流，不是刚性规则引擎。"],
        fill=LIGHT_ORANGE,
        line=ORANGE,
    )
    _add_textbox(slide, 7.05, 6.20, 5.4, 0.55, "阶段判断：项目已具备可用工程基础，下一阶段重点应转向业务场景深化与仿真接入。", size=14, bold=True, color=PURPLE_DARK)


def build_problem_slide(prs: Presentation):
    slide = _new_slide(prs, "项目目标与问题定义", "第 2 部分 / 设计起点")
    _add_card(
        slide,
        0.7,
        1.45,
        3.8,
        2.0,
        "问题 01：规范复杂",
        ["人工查阅 XSD、标准文档和元素关系成本高。", "难以快速定位父子关系、choice/sequence 与引用语义。"],
        fill=LIGHT_RED,
        line=RED,
        title_color=RED,
    )
    _add_card(
        slide,
        4.75,
        1.45,
        3.8,
        2.0,
        "问题 02：模型易错",
        ["大模型能写 XML，但在结构、变体、枚举、顺序上容易出错。", "复杂场景下容易把局部错误扩大成全局重写。"],
        fill=LIGHT_ORANGE,
        line=ORANGE,
        title_color=ORANGE,
    )
    _add_card(
        slide,
        8.8,
        1.45,
        3.8,
        2.0,
        "问题 03：纯规则引擎不够灵活",
        ["开放式场景需求下，模板拼接的泛化能力不足。", "复杂语义往往仍需大模型负责整体决策。"],
        fill=LIGHT_BLUE,
        line=BLUE,
        title_color=BLUE,
    )
    _add_flow_row(
        slide,
        0.9,
        4.0,
        11.3,
        0.65,
        ["自然语言需求", "结构化知识图谱", "MCP 结构辅助", "本地校验/修复", "VTD 场景输出"],
        fills=[LIGHT, LIGHT_BLUE, LIGHT_TEAL, LIGHT_ORANGE, LIGHT_RED],
        text_size=13,
    )
    _add_bullets(
        slide,
        1.05,
        4.95,
        11.0,
        1.65,
        [
            "目标不是替代大模型，而是在其自由决策外侧加上结构化护栏。",
            "最终方案选型为：LLM-first, MCP-assisted。",
            "MCP 只做高价值辅助：检索、schema 提示、验证、诊断、guidance 组合。",
        ],
        size=15,
        fill=WHITE,
        border=BORDER,
        title="技术路线选择",
    )


def build_architecture_slide(prs: Presentation):
    slide = _new_slide(prs, "总体技术架构", "第 3 部分 / 系统全景")
    _add_flow_row(
        slide,
        0.7,
        1.45,
        11.6,
        0.72,
        ["本地 XSD / 文档", "XSD Parser", "结构化知识图谱", "FastMCP Runtime", "Skill / Agent", "VTD 导向 XML"],
        fills=[LIGHT_ORANGE, LIGHT_BLUE, LIGHT, LIGHT_TEAL, LIGHT_BLUE, LIGHT_RED],
        text_size=13,
    )
    _add_card(
        slide,
        0.75,
        2.45,
        3.85,
        2.85,
        "原始知识层",
        [
            "schema/OpenSCENARIO.xsd",
            "raw/docs/",
            "raw/validator/",
            "作用：提供规范原始来源，不直接喂给模型决策。",
        ],
        fill=WHITE,
        line=PURPLE,
    )
    _add_card(
        slide,
        4.78,
        2.45,
        3.2,
        2.85,
        "结构化知识层",
        [
            "elements/*.json",
            "schema_scope.json",
            "coverage_report.json",
            "作用：把 XSD 转成可计算知识图谱。",
        ],
        fill=LIGHT,
        line=PURPLE,
    )
    _add_card(
        slide,
        8.12,
        2.45,
        4.15,
        2.85,
        "MCP / Workflow 层",
        [
            "retrieve_spec",
            "get_element_schema",
            "build_xml_guidance",
            "validate_xml / explain_validation_errors",
            "Skill 与 guidance packet 工作流",
        ],
        fill=LIGHT_BLUE,
        line=BLUE,
    )
    _add_bullets(
        slide,
        0.82,
        5.60,
        11.35,
        1.0,
        [
            "架构要点：知识静态沉淀，策略运行时推导，XML 仍由大模型负责最终设计与修复。",
        ],
        size=15,
        fill=LIGHT_TEAL,
        border=TEAL,
        title="评审结论",
    )


def build_principles_slide(prs: Presentation):
    slide = _new_slide(prs, "技术路线与设计原则", "第 4 部分 / 关键原则")
    _add_card(slide, 0.75, 1.55, 3.85, 3.15, "原则 01：模型拥有决策自由", ["不把大模型退化成模板填空器。", "保留其对复杂场景逻辑、元素选择和修复路径的主导权。"], fill=LIGHT_BLUE, line=BLUE)
    _add_card(slide, 4.75, 1.55, 3.05, 3.15, "原则 02：结构风险必须被机器约束", ["对 choice、sequence、contextual variant、reference、enum 做结构化约束。", "把高风险部分从“经验”升级为“可计算策略”。"], fill=LIGHT_TEAL, line=TEAL)
    _add_card(slide, 8.0, 1.55, 4.2, 3.15, "原则 03：错误修复必须可解释可追踪", ["统一进入 validate -> diagnose -> repair 的闭环。", "错误需定位、分类、给出优先修复动作，而不是靠人工猜。"], fill=LIGHT_ORANGE, line=ORANGE)
    _add_table(
        slide,
        0.9,
        5.15,
        11.2,
        1.25,
        [
            ["取舍问题", "本项目选择"],
            ["谁负责 XML 设计", "大模型负责"],
            ["谁负责结构约束", "知识图谱 + MCP"],
            ["谁负责最终语法闭环", "本地 Python validator"],
        ],
        [5.0, 6.2],
        header_fill=PURPLE,
        row_fill=WHITE,
    )


def build_kg_goal_slide(prs: Presentation):
    slide = _new_slide(prs, "知识图谱建设目标", "第 5 部分 / 知识层")
    _add_card(slide, 0.75, 1.55, 3.7, 2.4, "传统资料形态", ["文本分散", "语义难计算", "同名元素上下文难分辨", "难直接驱动模型修复"], fill=LIGHT_RED, line=RED, title_color=RED)
    _add_card(slide, 4.65, 1.55, 3.1, 2.4, "本项目改造方向", ["把元素作为节点", "把属性/顺序/变体/引用作为结构关系", "把 XSD 变成机器可消费 JSON"], fill=LIGHT_BLUE, line=BLUE, title_color=BLUE)
    _add_card(slide, 7.95, 1.55, 4.2, 2.4, "目标输出", ["可检索", "可追溯", "可推导 strategy", "可直接服务 XML 草拟与修复"], fill=LIGHT_TEAL, line=TEAL, title_color=TEAL)
    _add_flow_row(slide, 1.0, 4.35, 10.8, 0.68, ["元素节点", "结构关系", "上下文变体", "引用语义", "运行时策略"], fills=[LIGHT, LIGHT_BLUE, LIGHT_BLUE, LIGHT_ORANGE, LIGHT_TEAL], text_size=13)
    _add_bullets(slide, 1.0, 5.35, 11.0, 1.05, ["这里的“知识图谱”不是向量检索文档库，而是由元素记录、关系字段和运行时策略共同组成的结构化知识底座。"], size=15, fill=WHITE, border=BORDER, title="评审提示")


def build_kg_pipeline_slide(prs: Presentation):
    slide = _new_slide(prs, "知识图谱构建流水线", "第 6 部分 / 自动化生产链")
    _add_flow_row(slide, 0.75, 1.45, 11.55, 0.62, ["OpenSCENARIO.xsd", "xsd_inventory.py", "xsd_parser.py", "elements/*.json", "schema_scope.json", "coverage_report.json"], fills=[LIGHT_ORANGE, LIGHT_BLUE, LIGHT_BLUE, LIGHT, LIGHT_TEAL, LIGHT_RED], text_size=12)
    _add_card(slide, 0.8, 2.35, 3.8, 3.35, "阶段 01：清点与解析", ["自动统计 element / simpleType / complexType / group。", "逐元素解析 required attributes、allowed children、child order、enum。", "解析 extension、choice group、simpleContent 等结构细节。"])
    _add_card(slide, 4.8, 2.35, 3.25, 3.35, "阶段 02：语义补强", ["推断 parent_contexts。", "推断 contextual_variants。", "识别 reference_kind。", "生成 semantic_constraints。"], fill=LIGHT_BLUE, line=BLUE)
    _add_card(slide, 8.25, 2.35, 3.95, 3.35, "阶段 03：质量闭环", ["生成 schema_scope.json。", "生成 coverage_report.json。", "检查 missing elements / extra elements / dangling refs / metadata 缺失。"], fill=LIGHT_TEAL, line=TEAL)
    _add_textbox(slide, 0.88, 5.95, 11.1, 0.45, "结论：这是一条“从规范到结构资产”的自动生产链，而不是人工维护的静态词条库。", size=15, bold=True, color=PURPLE_DARK)


def build_kg_coverage_slide(prs: Presentation):
    slide = _new_slide(prs, "知识图谱覆盖规模与域划分", "第 7 部分 / 覆盖状态")
    _add_metric_card(slide, 0.85, 1.45, 2.45, 1.05, "302", "XSD 元素名总数")
    _add_metric_card(slide, 3.55, 1.45, 2.45, 1.05, "301", "结构化 JSON 文件数")
    _add_metric_card(slide, 6.25, 1.45, 2.45, 1.05, "0", "覆盖缺口")
    _add_metric_card(slide, 8.95, 1.45, 3.00, 1.05, "0 / 0", "悬挂引用 / 元数据缺失")
    _add_table(
        slide,
        0.85,
        2.85,
        7.15,
        2.95,
        [
            ["领域", "元素数量"],
            ["domain-core-entities", "77"],
            ["domain-routing-geometry", "55"],
            ["domain-actions-control", "72"],
            ["domain-conditions-values", "42"],
            ["domain-traffic-environment", "55"],
        ],
        [5.4, 1.6],
    )
    _add_card(slide, 8.3, 2.85, 3.8, 2.95, "覆盖状态补充说明", ["别名冲突单独追踪：OpenScenario -> OpenSCENARIO", "represented_xsd_elements = 302", "missing_elements = []", "extra_structured_elements = []"], fill=LIGHT_ORANGE, line=ORANGE)
    _add_textbox(slide, 0.92, 6.05, 11.1, 0.42, "评审结论：当前不是 MVP 子集，而是已经实现对本地 XSD 的全量结构化覆盖。", size=15, bold=True, color=PURPLE_DARK)


def build_kg_fields_slide(prs: Presentation):
    slide = _new_slide(prs, "元素记录字段设计", "第 8 部分 / 记录模型")
    _add_card(slide, 0.75, 1.45, 2.8, 2.7, "身份与上下文", ["element", "description", "parent_contexts", "contextual_variants"], fill=LIGHT_BLUE, line=BLUE)
    _add_card(slide, 3.7, 1.45, 2.8, 2.7, "结构字段", ["required_attributes", "optional_attributes", "allowed_children", "child_order", "multiplicity"], fill=WHITE, line=PURPLE)
    _add_card(slide, 6.65, 1.45, 2.8, 2.7, "语义字段", ["content_model_kind", "child_groups", "semantic_constraints", "enum_constraints", "reference_kind"], fill=LIGHT_TEAL, line=TEAL)
    _add_card(slide, 9.6, 1.45, 2.2, 2.7, "追溯字段", ["source_path"], fill=LIGHT_ORANGE, line=ORANGE)
    _add_code_block(
        slide,
        0.85,
        4.45,
        5.75,
        1.85,
        [
            '{',
            '  "element": "Storyboard",',
            '  "parent_contexts": ["ScenarioDefinition"],',
            '  "allowed_children": ["Init", "Story", "StopTrigger"],',
            '  "child_order": ["Init", "Story", "StopTrigger"]',
            '}',
        ],
        title="结构化记录示例",
    )
    _add_bullets(slide, 6.85, 4.45, 5.1, 1.95, ["字段设计目标不是“存资料”，而是让模型能直接拿到结构、限制和来源。", "这些字段会在运行时进一步生成 strategy，用于草拟与修复。"], size=15, fill=WHITE, border=BORDER, title="设计意义")


def build_action_example_slide(prs: Presentation):
    slide = _new_slide(prs, "实例 01：Action 节点", "第 9 部分 / choice 结构示例")
    _add_code_block(
        slide,
        0.75,
        1.45,
        5.65,
        3.95,
        [
            '{',
            '  "element": "Action",',
            '  "content_model_kind": "choice",',
            '  "child_groups": [{"members":',
            '    ["GlobalAction","UserDefinedAction","PrivateAction"],',
            '    "cardinality": "1..1"}],',
            '  "required_attributes": [{"name": "name"}],',
            '  "parent_contexts": ["Event"]',
            '}',
        ],
        title="Action.json 摘要",
    )
    _add_card(slide, 6.7, 1.45, 5.45, 2.1, "这条记录实际告诉模型什么", ["Action 不是普通容器，而是 choice 包装节点。", "GlobalAction / UserDefinedAction / PrivateAction 必须三选一。", "该元素只能处于 Event 上下文中。", "还必须带 name 属性。"], fill=LIGHT_BLUE, line=BLUE)
    _add_card(slide, 6.7, 3.85, 5.45, 1.55, "为什么重要", ["如果没有 child_groups 和 content_model_kind，模型极易把多个分支一起塞进 Action。"], fill=LIGHT_ORANGE, line=ORANGE)
    _add_flow_row(slide, 1.1, 5.85, 10.8, 0.62, ["Action", "GlobalAction | UserDefinedAction | PrivateAction", "输出单一合法分支"], fills=[LIGHT, LIGHT_BLUE, LIGHT_ORANGE], text_size=13)


def build_catalog_example_slide(prs: Presentation):
    slide = _new_slide(prs, "实例 02：CatalogReference 变体节点", "第 10 部分 / contextual variants 示例")
    _add_card(
        slide,
        0.75,
        1.45,
        5.85,
        4.0,
        "CatalogReference 的关键特征",
        [
            "必填属性：catalogName、entryName",
            "可选子节点：ParameterAssignments",
            "子节点顺序：ParameterAssignments",
            "可出现在多个父上下文：AssignRouteAction、RouteRef、TrajectoryRef、EnvironmentAction、ObjectController 等",
        ],
        fill=WHITE,
        line=PURPLE,
    )
    _add_table(
        slide,
        6.85,
        1.45,
        5.05,
        4.0,
        [
            ["典型父上下文", "含义"],
            ["AssignRouteAction", "路由分配路径"],
            ["FollowTrajectoryAction", "轨迹跟随路径"],
            ["EnvironmentAction", "环境目录引用"],
            ["RouteRef / TrajectoryRef", "引用包装节点"],
            ["ObjectController", "控制器相关目录引用"],
        ],
        [2.45, 2.45],
        header_fill=PURPLE,
        row_fill=WHITE,
    )
    _add_bullets(slide, 0.85, 5.8, 11.1, 0.62, ["评审意义：模型不是只要知道“CatalogReference 长什么样”，而是要知道“当前应挂在哪个父节点下”。"], size=15, fill=LIGHT_TEAL, border=TEAL)


def build_kg_drives_slide(prs: Presentation):
    slide = _new_slide(prs, "知识图谱如何驱动大模型决策", "第 11 部分 / 从静态知识到动态作用")
    _add_table(
        slide,
        0.8,
        1.45,
        11.35,
        3.15,
        [
            ["风险类型", "对应字段", "给模型的直接帮助"],
            ["choice 误写", "content_model_kind / child_groups", "知道是三选一还是可多选"],
            ["顺序错误", "child_order", "知道 sequence 中谁必须在前"],
            ["同名元素误用", "contextual_variants / parent_contexts", "结合父上下文分辨合法变体"],
            ["引用漏接", "reference_kind", "知道应接变量/实体/轨迹等具体引用"],
            ["枚举乱填", "enum_constraints", "知道哪些字面值才合法"],
        ],
        [2.2, 3.4, 5.4],
        header_fill=PURPLE,
        row_fill=WHITE,
    )
    _add_flow_row(slide, 1.0, 5.15, 10.8, 0.62, ["用户意图", "遇到结构风险", "查询图谱记录", "读取策略字段", "生成更稳健 XML"], fills=[LIGHT, LIGHT_ORANGE, LIGHT_BLUE, LIGHT_TEAL, LIGHT_RED], text_size=13)
    _add_textbox(slide, 1.0, 6.0, 11.0, 0.4, "结论：知识图谱不是“多一些说明文本”，而是在风险节点上直接参与模型的结构决策。", size=15, bold=True, color=PURPLE_DARK)


def build_mcp_overview_slide(prs: Presentation):
    slide = _new_slide(prs, "MCP 服务总览", "第 12 部分 / Runtime 全景")
    _add_card(slide, 0.75, 1.45, 3.3, 2.4, "Runtime 内核", ["KnowledgeBase", "diagnostic patterns", "ValidatorAdapter", "FastMCP server"], fill=WHITE, line=PURPLE)
    _add_card(slide, 4.2, 1.45, 3.9, 2.4, "结构智能层", ["retrieve_spec", "get_element_schema", "build_element_strategy", "search / source trace"], fill=LIGHT_BLUE, line=BLUE)
    _add_card(slide, 8.3, 1.45, 3.9, 2.4, "校验与修复层", ["validate_xml", "explain_validation_errors", "repair_strategy", "guidance packet"], fill=LIGHT_TEAL, line=TEAL)
    _add_flow_row(slide, 0.95, 4.25, 11.1, 0.62, ["retrieve_spec", "get_element_schema", "build_xml_guidance", "validate_xml", "explain_validation_errors"], fills=[LIGHT, LIGHT_BLUE, LIGHT_BLUE, LIGHT_ORANGE, LIGHT_RED], text_size=12)
    _add_bullets(slide, 0.95, 5.15, 11.05, 1.25, ["MCP 不直接生成 XML，而是向模型提供有边界、可解释、可追溯的结构辅助。", "技术亮点在于：静态 schema 记录在运行时被进一步转成 strategy，形成“搜索 + 结构 + 诊断”的组合能力。"], size=15, fill=WHITE, border=BORDER, title="评审要点")


def build_retrieve_slide(prs: Presentation):
    slide = _new_slide(prs, "模块详解：retrieve_spec", "第 13 部分 / 检索模块")
    _add_flow_row(slide, 0.75, 1.45, 7.0, 0.62, ["自然语言 query", "文本规范化", "候选集构建", "打分排序", "返回结构摘要"], fills=[LIGHT, LIGHT_BLUE, LIGHT_BLUE, LIGHT_TEAL, LIGHT_ORANGE], text_size=12)
    _add_card(slide, 0.78, 2.35, 6.95, 2.05, "已完成能力", ["支持 element / attribute / concept / error 四类检索。", "支持 CamelCase 名称与自然语言混合匹配。", "支持 parent_context 参与结果判别。", "返回 source_path、constraints、strategy_summary。"])
    _add_bullets(slide, 8.0, 1.45, 4.15, 2.9, ["检索不是只找资料，而是为“当前这一步 XML 编写动作”筛选最相关结构信息。", "strategy_summary 会提前告诉模型：是否存在 choice、顺序约束、上下文变体或必填引用。"], size=14, fill=LIGHT_BLUE, border=BLUE, title="模块价值")
    _add_code_block(slide, 8.0, 4.6, 4.15, 1.65, ['query = "Storyboard"', 'kind = "element"', "", "=> Required children: Init", "=> Child order: Init -> Story -> StopTrigger", "=> strategy_summary: preserve order, keep Init present"], title="真实返回示意")


def build_schema_slide(prs: Presentation):
    slide = _new_slide(prs, "模块详解：get_element_schema", "第 14 部分 / schema 查询")
    _add_flow_row(slide, 0.75, 1.45, 7.1, 0.62, ["element + parent_context", "取 ElementRecord", "build_element_strategy", "返回 record + strategy"], fills=[LIGHT, LIGHT_BLUE, LIGHT_TEAL, LIGHT_ORANGE], text_size=12)
    _add_card(slide, 0.8, 2.35, 7.0, 2.2, "strategy 关键字段", ["structure_mode", "branch_selection", "ordering", "required_children", "variant_resolution", "reference_requirements", "repair_priority"], fill=WHITE, line=PURPLE)
    _add_code_block(slide, 8.05, 1.45, 4.0, 2.1, ['element = "SetAction"', 'parent_context = "VariableAction"', "", 'resolved_variant = {', '  "type_ref": "VariableSetAction",', '  "deprecated": false', '}'], title="上下文分辨示例")
    _add_bullets(slide, 8.05, 3.8, 4.0, 2.45, ["这是把静态 schema 升级成 machine-readable strategy 的关键模块。", "同名元素会结合 parent_context 做非 deprecated 优先分辨。", "缺失 variableRef 这类问题会被识别成“先补引用”，而不是泛化为“补属性”。"], size=14, fill=LIGHT_TEAL, border=TEAL, title="评审意义")


def build_validate_slide(prs: Presentation):
    slide = _new_slide(prs, "模块详解：validate_xml", "第 15 部分 / 本地校验")
    _add_flow_row(slide, 0.75, 1.45, 7.0, 0.62, ["XML 字符串", "UTF-8 解析", "XMLSchema.validate", "归一化错误对象"], fills=[LIGHT, LIGHT_BLUE, LIGHT_ORANGE, LIGHT_RED], text_size=12)
    _add_card(slide, 0.82, 2.35, 6.95, 2.1, "已完成能力", ["基于 lxml.etree.XMLSchema 编译本地 OpenSCENARIO.xsd。", "支持 1.4.0 / 1.4 / 1.x / v1.4.0 / v1.4 版本别名。", "统一输出 line / column / message / rule_hint。"], fill=WHITE, line=PURPLE)
    _add_code_block(slide, 8.0, 1.45, 4.1, 2.3, ["<ManeuverGroup />", "", "=> ok = false", "=> message = Element 'ManeuverGroup': Missing", "   child element(s). Expected is ( Actors ).", '=> rule_hint = "Actors"'], title="真实校验示例")
    _add_bullets(slide, 8.0, 4.0, 4.1, 1.8, ["这是整个系统最确定性的环节。", "没有它，前面所有结构辅助都无法形成最终闭环。"], size=15, fill=LIGHT_ORANGE, border=ORANGE, title="评审意义")


def build_diagnostics_slide(prs: Presentation):
    slide = _new_slide(prs, "模块详解：explain_validation_errors", "第 16 部分 / 诊断模块")
    _add_flow_row(slide, 0.75, 1.45, 7.0, 0.62, ["errors", "regex 分类", "schema-aware enrich", "repair_strategy"], fills=[LIGHT, LIGHT_BLUE, LIGHT_TEAL, LIGHT_ORANGE], text_size=12)
    _add_card(slide, 0.82, 2.35, 6.95, 2.25, "已覆盖的主要类别", ["namespace_or_root_issue", "missing_required_child", "missing_required_attribute", "wrong_child_order", "invalid_attribute", "invalid_enum_value", "unexpected_element"], fill=WHITE, line=PURPLE)
    _add_code_block(slide, 8.0, 1.45, 4.1, 2.1, ["Element 'VariableAction':", "The attribute 'variableRef' is required", "but missing.", "", '=> recommended_actions = ["add_required_references"]'], title="引用感知修复示例")
    _add_bullets(slide, 8.0, 3.8, 4.1, 2.2, ["诊断模块不是简单翻译器，而是会结合知识图谱判断：", "是否是 choice 缺分支、是否只是顺序错误、是否应该先补引用。", "最终返回 focus_element、expected_elements 和 recommended_actions。"], size=14, fill=LIGHT_RED, border=RED, title="模块价值")


def build_guidance_slide(prs: Presentation):
    slide = _new_slide(prs, "模块详解：build_xml_guidance", "第 17 部分 / guidance 组合层")
    _add_flow_row(slide, 0.75, 1.45, 7.1, 0.62, ["query + element", "retrieve_spec", "get_element_schema", "errors? -> explain", "guidance packet"], fills=[LIGHT, LIGHT_BLUE, LIGHT_BLUE, LIGHT_ORANGE, LIGHT_TEAL], text_size=12)
    _add_card(slide, 0.82, 2.35, 6.95, 2.15, "packet 内容", ["retrieval_hits", "element_schema", "draft_checklist", "repair_diagnostics", "repair_actions"], fill=WHITE, line=PURPLE)
    _add_code_block(slide, 8.0, 1.45, 4.1, 2.4, ['query = "set action"', 'element = "SetAction"', 'parent_context = "VariableAction"', "", "=> draft_checklist:", "   1. 选 VariableSetAction 变体", "   2. Value / Expression 二选一", "   3. 先接 variableRef"], title="真实 guidance 片段")
    _add_bullets(slide, 8.0, 4.1, 4.1, 1.95, ["组合层的意义是降低 agent 拼装上下文成本。", "尤其适合新会话、外部 agent 和 benchmark guidance packet 生成。"], size=15, fill=LIGHT_TEAL, border=TEAL, title="评审结论")


def build_tool_loop_slide(prs: Presentation):
    slide = _new_slide(prs, "MCP 五工具协同闭环", "第 18 部分 / 工具协作")
    lane_y = 1.55
    names = ["用户/Agent", "retrieve", "schema", "guidance", "validate", "diagnostics"]
    x_positions = [0.8, 2.7, 4.75, 6.8, 9.05, 11.1]
    fills = [LIGHT, LIGHT_BLUE, LIGHT_BLUE, LIGHT_TEAL, LIGHT_ORANGE, LIGHT_RED]
    for x, name, fill in zip(x_positions, names, fills):
        _add_card(slide, x, lane_y, 1.55, 0.75, name, [], fill=fill, line=PURPLE if name == "用户/Agent" else BLUE, title_size=11)
        _add_connector(slide, x + 0.78, 2.3, x + 0.78, 5.8, color=BORDER, width=1.0)
    steps = [
        ("1. 查询行为/元素", 1.15, 2.55, 3.4, 0.45, LIGHT),
        ("2. 返回 hits + 结构摘要", 2.95, 3.00, 3.2, 0.45, LIGHT_BLUE),
        ("3. 获取完整 schema + strategy", 4.9, 3.55, 3.15, 0.45, LIGHT_BLUE),
        ("4. 需要时组合 guidance packet", 6.9, 4.10, 3.0, 0.45, LIGHT_TEAL),
        ("5. 提交 XML 校验", 8.95, 4.65, 2.25, 0.45, LIGHT_ORANGE),
        ("6. 校验失败 -> 诊断修复动作", 9.75, 5.20, 2.55, 0.45, LIGHT_RED),
    ]
    for text, left, top, width, height, fill in steps:
        _add_card(slide, left, top, width, height, text, [], fill=fill, line=BORDER, title_size=12, body_size=11)
    _add_textbox(slide, 0.95, 6.15, 11.0, 0.35, "闭环价值：先找知识，再看结构，再校验，再诊断，再局部修复，避免全局重写。", size=15, bold=True, color=PURPLE_DARK)


def build_skill_slide(prs: Presentation):
    slide = _new_slide(prs, "Skill 工作流", "第 19 部分 / LLM-first 编排")
    _add_flow_row(slide, 0.75, 1.45, 11.45, 0.62, ["场景需求", "parsed_intent", "xml_intent_check", "LLM 起草 XML", "validate_xml", "intent_consistent"], fills=[LIGHT, LIGHT_BLUE, LIGHT_BLUE, LIGHT_TEAL, LIGHT_ORANGE, LIGHT_RED], text_size=12)
    _add_card(slide, 0.82, 2.35, 3.75, 2.9, "Skill 固化的中间状态", ["parsed_intent", "xml_intent_check", "schema_valid", "intent_consistent", "remaining_blockers"], fill=WHITE, line=PURPLE)
    _add_card(slide, 4.78, 2.35, 3.35, 2.9, "何时调用 MCP", ["遇到结构风险时：retrieve_spec / get_element_schema", "已知目标元素时：build_xml_guidance", "校验失败时：validate_xml + explain_validation_errors"], fill=LIGHT_BLUE, line=BLUE)
    _add_card(slide, 8.33, 2.35, 3.9, 2.9, "VTD 导向约束", ["偏向保守、最小、可运行的结构。", "避免无根据的 catalog/controller 扩写。", "短语如“OpenX 场景代码”“VTD 场景文件”都会命中该工作流。"], fill=LIGHT_TEAL, line=TEAL)
    _add_textbox(slide, 0.95, 5.65, 11.0, 0.55, "结论：Skill 把“模型自由度”和“结构安全性”放进同一条执行链，而不是让二者互相冲突。", size=15, bold=True, color=PURPLE_DARK)


def build_external_workflow_slide(prs: Presentation):
    slide = _new_slide(prs, "外部 Agent / Guidance Packet 工作流", "第 20 部分 / 跨会话复用")
    _add_flow_row(slide, 0.8, 1.45, 11.35, 0.62, ["prompt 文件", "build_guidance_packet.py", ".guidance.json", "新 Codex / 外部 agent", "XML + 本地校验"], fills=[LIGHT, LIGHT_BLUE, LIGHT_TEAL, LIGHT_BLUE, LIGHT_ORANGE], text_size=12)
    _add_card(slide, 0.85, 2.35, 4.0, 2.6, "为什么需要 guidance packet", ["复杂任务跨会话时，模型上下文会重建。", "packet 把结构风险、候选元素、修复动作沉淀成一个可携带对象。"], fill=WHITE, line=PURPLE)
    _add_card(slide, 5.05, 2.35, 3.35, 2.6, "packet 的定位", ["它是 advisory，不替代模型决策。", "它帮助缩小结构判断空间，但最终 XML 仍由模型决定。"], fill=LIGHT_BLUE, line=BLUE)
    _add_card(slide, 8.6, 2.35, 3.55, 2.6, "对应文档与脚本", ["docs/external-agent-guidance-workflow.md", "docs/external-agent-prompt-template.md", "scripts/build_guidance_packet.py", "scripts/build_benchmark_guidance.py"], fill=LIGHT_TEAL, line=TEAL)
    _add_textbox(slide, 0.95, 5.55, 11.0, 0.65, "评审结论：该机制显著提高了跨会话交接稳定性，同时保持了 LLM-first 的总体路线。", size=15, bold=True, color=PURPLE_DARK)


def build_benchmark_slide(prs: Presentation):
    slide = _new_slide(prs, "Benchmark 与 Sidecar 契约", "第 21 部分 / 可复现验证资产")
    _add_card(slide, 0.75, 1.45, 4.05, 3.2, "当前 benchmark 集", ["minimal-single-vehicle", "two-vehicle-follow", "triggered-deceleration", "triggered-lane-change"], fill=WHITE, line=PURPLE)
    _add_card(slide, 5.0, 1.45, 3.15, 3.2, "guidance-inputs.json", ["每个 benchmark 绑定 query、target element、parent_context。", "用于自动生成 .guidance.json。"], fill=LIGHT_BLUE, line=BLUE)
    _add_card(slide, 8.35, 1.45, 3.8, 3.2, "intent-schema.json", ["统一 sidecar 字段：", "parsed_intent", "xml_intent_check", "schema_valid", "intent_consistent", "remaining_blockers"], fill=LIGHT_TEAL, line=TEAL)
    _add_bullets(slide, 0.95, 5.05, 11.0, 1.2, ["benchmark 不是只存 XML 结果，而是把提示词、意图、检查单和校验状态一起纳入回归对象。", "run-log.md 当前记录 4 个 benchmark 均为 pass。"], size=15, fill=LIGHT_ORANGE, border=ORANGE, title="评审要点")


def build_verification_slide(prs: Presentation):
    slide = _new_slide(prs, "验证体系与当前结果", "第 22 部分 / 工程质量")
    _add_card(slide, 0.8, 1.45, 3.7, 3.35, "覆盖维度 01", ["xsd inventory", "schema coverage report", "full schema scope", "domain knowledge correctness"], fill=WHITE, line=PURPLE)
    _add_card(slide, 4.7, 1.45, 3.5, 3.35, "覆盖维度 02", ["generation strategy", "schema tool", "retrieve_spec tool", "diagnostics tool", "guidance tool / runner"], fill=LIGHT_BLUE, line=BLUE)
    _add_card(slide, 8.4, 1.45, 3.75, 3.35, "覆盖维度 03", ["server registration", "tool loop integration", "benchmark assets", "benchmark results"], fill=LIGHT_TEAL, line=TEAL)
    _add_metric_card(slide, 1.25, 5.2, 2.5, 1.0, "93 passed", "最新完整回归")
    _add_metric_card(slide, 4.35, 5.2, 2.5, 1.0, "0", "缺失元素")
    _add_metric_card(slide, 7.45, 5.2, 2.5, 1.0, "0", "悬挂引用")
    _add_metric_card(slide, 10.25, 5.2, 1.75, 1.0, "5", "工具数")
    _add_textbox(slide, 0.9, 6.35, 11.1, 0.32, "结论：当前交付物不是概念演示，而是具备自动化回归能力的工程化系统。", size=15, bold=True, color=PURPLE_DARK)


def build_final_slide(prs: Presentation):
    slide = _new_slide(prs, "当前阶段交付与下一步方向", "第 23 部分 / 收口")
    _add_card(slide, 0.75, 1.45, 5.6, 4.4, "当前阶段已交付", ["1. 全量结构化知识图谱底座", "2. FastMCP 服务与 5 个核心工具", "3. 生成 / 修复 strategy 层", "4. Codex Skill 与外部 agent 工作流", "5. benchmark 资产、sidecar 契约与自动化回归", "6. 中文使用文档与阶段汇报材料"], fill=WHITE, line=PURPLE)
    _add_card(slide, 6.6, 1.45, 5.55, 2.05, "阶段判断", ["项目已完成从“文档+人工查 schema”到“知识图谱 + MCP + 校验闭环”的跃迁。", "下一阶段重点不再是“能不能生成 XML”，而是“如何更深贴合 VTD 业务场景与仿真执行链”。"], fill=LIGHT_BLUE, line=BLUE)
    _add_card(slide, 6.6, 3.75, 5.55, 2.1, "建议的下一步", ["扩展真实业务 benchmark 与场景族。", "沉淀更强的语义知识与业务约束。", "推进与仿真执行、资产管理、结果回放链路的深度集成。"], fill=LIGHT_TEAL, line=TEAL)
    _add_textbox(slide, 0.95, 6.15, 11.1, 0.35, "总体结论：这套系统已经具备继续向甲方业务深水区推进的坚实工程底座。", size=16, bold=True, color=PURPLE_DARK)


SLIDES = [
    build_cover,
    build_scope_slide,
    build_problem_slide,
    build_architecture_slide,
    build_principles_slide,
    build_kg_goal_slide,
    build_kg_pipeline_slide,
    build_kg_coverage_slide,
    build_kg_fields_slide,
    build_action_example_slide,
    build_catalog_example_slide,
    build_kg_drives_slide,
    build_mcp_overview_slide,
    build_retrieve_slide,
    build_schema_slide,
    build_validate_slide,
    build_diagnostics_slide,
    build_guidance_slide,
    build_tool_loop_slide,
    build_skill_slide,
    build_external_workflow_slide,
    build_benchmark_slide,
    build_verification_slide,
    build_final_slide,
]


def build_deck():
    template = _template_path()
    output = _output_path()
    prs = Presentation(str(template))
    _clear_all_slides(prs)
    for builder in SLIDES:
        builder(prs)
    prs.save(str(output))
    try:
        assets_dir = _template_path().parent / "assets" / "mermaid"
        if assets_dir.is_dir():
            from apply_report_mermaid_diagrams import apply_diagram_images

            apply_diagram_images(output, assets_dir)
    except Exception as exc:  # pragma: no cover
        print(f"Skipped Mermaid diagram overlay: {exc}")
    print(f"Wrote PPT to {output}")


if __name__ == "__main__":
    build_deck()
